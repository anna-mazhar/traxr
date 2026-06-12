"""PowerPoint tool for structured PPTX access."""

import base64
import io
from pathlib import Path
from typing import Any, Dict, List, Optional
from .base import BaseTool, ToolResult

# Lazy imports
HAS_PPTX = None
HAS_PIL = None


def _check_pptx():
    """Check if python-pptx is available."""
    global HAS_PPTX
    if HAS_PPTX is None:
        try:
            from pptx import Presentation
            HAS_PPTX = True
        except ImportError:
            HAS_PPTX = False
    return HAS_PPTX


def _check_pil():
    """Check if PIL is available."""
    global HAS_PIL
    if HAS_PIL is None:
        try:
            from PIL import Image
            HAS_PIL = True
        except ImportError:
            HAS_PIL = False
    return HAS_PIL


class PPTXTool(BaseTool):
    """Tool for PowerPoint file operations with structured access.

    Provides operations:
    - load: Load PPTX and extract content
    - get_text: Get all text from slides
    - get_slide_count: Get number of slides
    - get_slide: Get content of a specific slide
    - get_images: Get images from slides (as base64)
    - get_summary: Get structural summary of the presentation
    """

    def __init__(self, file_path: Optional[str] = None):
        """Initialize PPTXTool.

        Args:
            file_path: Path to PPTX file (can be set later via load)
        """
        super().__init__(name="pptx")
        self.file_path = file_path
        self.presentation = None
        self.slides_content: List[Dict[str, Any]] = []
        self.images: List[Dict[str, Any]] = []
        self.metadata: Dict[str, Any] = {}
        # Perturbation support: injected content overrides extracted text
        self._injected_content: Optional[str] = None
        self._injection_active: bool = False

    def inject_perturbed_content(self, content: str) -> None:
        """Inject perturbed content to override PPTX text extraction.

        This allows perturbation experiments to modify the extracted text
        while keeping the same file type and processing path.

        Args:
            content: Perturbed text content to use instead of actual PPTX text
        """
        self._injected_content = content
        self._injection_active = True

    def clear_injection(self) -> None:
        """Clear any injected content, returning to normal PPTX extraction."""
        self._injected_content = None
        self._injection_active = False

    @property
    def has_injected_content(self) -> bool:
        """Check if perturbed content has been injected."""
        return self._injection_active and self._injected_content is not None

    def execute(self, operation: str, **kwargs) -> ToolResult:
        """Execute PPTX operation."""
        operations = {
            "load": self._load,
            "get_text": self._get_text,
            "get_slide_count": self._get_slide_count,
            "get_slide": self._get_slide,
            "get_images": self._get_images,
            "get_summary": self._get_summary,
        }

        if operation not in operations:
            return ToolResult(
                success=False,
                output=None,
                error=f"Unknown operation '{operation}'. Available: {list(operations.keys())}"
            )

        try:
            return operations[operation](**kwargs)
        except Exception as e:
            return ToolResult(
                success=False,
                output=None,
                error=f"Error in {operation}: {str(e)}"
            )

    def get_available_operations(self) -> List[str]:
        """Get list of available operations."""
        return ["load", "get_text", "get_slide_count", "get_slide", "get_images", "get_summary"]

    def get_schema(self) -> "ToolSchema":
        from .tool_schema import ToolSchema, OperationSchema, ToolParameterSchema
        return ToolSchema(
            name="pptx",
            description="PowerPoint tool for loading PPTX files, extracting text, images, and slide content.",
            operations={
                "load": OperationSchema(
                    name="load",
                    description="Load a PPTX file and extract all content (text, tables, images).",
                    parameters=[
                        ToolParameterSchema(name="file_path", type="string", description="Path to PPTX file", required=False),
                    ],
                ),
                "get_text": OperationSchema(
                    name="get_text",
                    description="Get all text from every slide in the presentation.",
                    parameters=[],
                ),
                "get_slide_count": OperationSchema(
                    name="get_slide_count",
                    description="Get the number of slides.",
                    parameters=[],
                ),
                "get_slide": OperationSchema(
                    name="get_slide",
                    description="Get content of a specific slide (text, shapes, tables).",
                    parameters=[
                        ToolParameterSchema(name="slide_number", type="integer", description="Slide number (1-based)", required=False, default=1),
                    ],
                ),
                "get_images": OperationSchema(
                    name="get_images",
                    description="Get images from slides as base64-encoded data.",
                    parameters=[
                        ToolParameterSchema(name="slide_number", type="integer", description="Filter to a specific slide number (1-based)", required=False),
                    ],
                ),
                "get_summary": OperationSchema(
                    name="get_summary",
                    description="Get structural summary of the presentation (slide count, images, text previews).",
                    parameters=[],
                ),
            },
        )

    def _load(self, file_path: Optional[str] = None) -> ToolResult:
        """Load PPTX file and extract content.

        Args:
            file_path: Path to PPTX file (optional if set in __init__)

        Returns:
            ToolResult with presentation info
        """
        if not _check_pptx():
            return ToolResult(
                success=False,
                output=None,
                error="python-pptx not installed. Install with: pip install python-pptx"
            )

        path = file_path or self.file_path
        if not path:
            return ToolResult(
                success=False,
                output=None,
                error="No file path provided"
            )

        path = Path(path)
        if not path.exists():
            return ToolResult(
                success=False,
                output=None,
                error=f"File not found: {path}"
            )

        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt

            self.file_path = str(path)
            self.presentation = Presentation(path)

            # Extract content from each slide
            self.slides_content = []
            self.images = []

            for slide_idx, slide in enumerate(self.presentation.slides):
                slide_data = {
                    "slide_number": slide_idx + 1,
                    "texts": [],
                    "shapes": [],
                    "has_images": False,
                }

                for shape in slide.shapes:
                    shape_info = {
                        "type": shape.shape_type.name if hasattr(shape.shape_type, 'name') else str(shape.shape_type),
                        "name": shape.name,
                    }

                    # Extract text
                    if shape.has_text_frame:
                        text_content = []
                        for paragraph in shape.text_frame.paragraphs:
                            para_text = paragraph.text.strip()
                            if para_text:
                                text_content.append(para_text)
                        if text_content:
                            slide_data["texts"].extend(text_content)
                            shape_info["text"] = "\n".join(text_content)

                    # Extract tables
                    if shape.has_table:
                        table = shape.table
                        table_data = []
                        for row in table.rows:
                            row_data = []
                            for cell in row.cells:
                                row_data.append(cell.text.strip())
                            table_data.append(row_data)
                        shape_info["table"] = table_data
                        # Also add table content as text
                        for row in table_data:
                            slide_data["texts"].append(" | ".join(row))

                    # Check for images
                    if shape.shape_type.name == "PICTURE" if hasattr(shape.shape_type, 'name') else False:
                        slide_data["has_images"] = True
                        try:
                            image = shape.image
                            image_data = {
                                "slide_number": slide_idx + 1,
                                "content_type": image.content_type,
                                "blob": base64.b64encode(image.blob).decode('utf-8'),
                                "filename": getattr(image, 'filename', f"slide_{slide_idx + 1}_image"),
                            }
                            self.images.append(image_data)
                        except Exception:
                            pass  # Skip images that can't be extracted

                    slide_data["shapes"].append(shape_info)

                self.slides_content.append(slide_data)

            # Store metadata
            self.metadata = {
                "file_name": path.name,
                "slide_count": len(self.slides_content),
                "total_images": len(self.images),
                "has_images": len(self.images) > 0,
            }

            return ToolResult(
                success=True,
                output=f"Loaded {path.name}: {len(self.slides_content)} slides, {len(self.images)} images",
                metadata=self.metadata
            )

        except Exception as e:
            return ToolResult(
                success=False,
                output=None,
                error=f"Failed to load PPTX: {str(e)}"
            )

    def _get_text(self) -> ToolResult:
        """Get all text from the presentation."""
        # Check for injected perturbed content (for perturbation experiments)
        if self._injection_active and self._injected_content:
            return ToolResult(
                success=True,
                output=self._injected_content,
                metadata={"slide_count": len(self.slides_content), "perturbed": True}
            )

        if not self.slides_content:
            return ToolResult(
                success=False,
                output=None,
                error="No presentation loaded. Call 'load' first."
            )

        all_text = []
        for slide in self.slides_content:
            slide_text = f"\n=== Slide {slide['slide_number']} ===\n"
            if slide["texts"]:
                slide_text += "\n".join(slide["texts"])
            else:
                slide_text += "(No text content)"
            all_text.append(slide_text)

        return ToolResult(
            success=True,
            output="\n".join(all_text),
            metadata={"slide_count": len(self.slides_content)}
        )

    def _get_slide_count(self) -> ToolResult:
        """Get number of slides."""
        if not self.slides_content:
            return ToolResult(
                success=False,
                output=None,
                error="No presentation loaded. Call 'load' first."
            )

        return ToolResult(
            success=True,
            output=len(self.slides_content),
            metadata=self.metadata
        )

    def _get_slide(self, slide_number: int = 1) -> ToolResult:
        """Get content of a specific slide.

        Args:
            slide_number: 1-based slide number

        Returns:
            ToolResult with slide content
        """
        if not self.slides_content:
            return ToolResult(
                success=False,
                output=None,
                error="No presentation loaded. Call 'load' first."
            )

        if slide_number < 1 or slide_number > len(self.slides_content):
            return ToolResult(
                success=False,
                output=None,
                error=f"Invalid slide number {slide_number}. Valid range: 1-{len(self.slides_content)}"
            )

        slide = self.slides_content[slide_number - 1]
        return ToolResult(
            success=True,
            output=slide,
            metadata={"slide_number": slide_number}
        )

    def _get_images(self, slide_number: Optional[int] = None) -> ToolResult:
        """Get images from slides.

        Args:
            slide_number: Optional 1-based slide number to filter images

        Returns:
            ToolResult with list of images (base64 encoded)
        """
        if not self.slides_content:
            return ToolResult(
                success=False,
                output=None,
                error="No presentation loaded. Call 'load' first."
            )

        if not self.images:
            return ToolResult(
                success=True,
                output=[],
                metadata={"message": "No images found in presentation"}
            )

        if slide_number:
            filtered = [img for img in self.images if img["slide_number"] == slide_number]
            return ToolResult(
                success=True,
                output=filtered,
                metadata={"image_count": len(filtered), "slide_number": slide_number}
            )

        return ToolResult(
            success=True,
            output=self.images,
            metadata={"image_count": len(self.images)}
        )

    def _get_summary(self) -> ToolResult:
        """Get structural summary of the presentation."""
        if not self.slides_content:
            return ToolResult(
                success=False,
                output=None,
                error="No presentation loaded. Call 'load' first."
            )

        summary = {
            "file_name": self.metadata.get("file_name", "unknown"),
            "slide_count": len(self.slides_content),
            "total_images": len(self.images),
            "slides": []
        }

        for slide in self.slides_content:
            slide_summary = {
                "slide_number": slide["slide_number"],
                "text_count": len(slide["texts"]),
                "has_images": slide["has_images"],
                "preview": slide["texts"][0][:100] if slide["texts"] else "(No text)",
            }
            summary["slides"].append(slide_summary)

        return ToolResult(
            success=True,
            output=summary,
            metadata=self.metadata
        )
